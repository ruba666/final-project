import ast

import matplotlib.pyplot as plt
import numpy as np
from openai import OpenAI
import io
import re
import os
from typing import List, Dict, Any
import random
from collections import defaultdict
from docx import Document as WordDocument
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.styles.differential import DifferentialStyle
from openpyxl.formatting.rule import Rule
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
import os
from typing import List, Dict, Any
import json
import random
from collections import defaultdict
from docx import Document as WordDocument
from docx.shared import Inches, Pt, RGBColor
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.styles.differential import DifferentialStyle
from openpyxl.formatting.rule import Rule
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
import re
import os
from typing import List, Dict, Any
import json
import random
from collections import defaultdict
from docx import Document as WordDocument
from docx.shared import Inches, Pt, RGBColor
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import matplotlib.pyplot as plt
import io


class CustomDocument:
    def __init__(self, name: str, doc_type: str):
        self.name = name
        self.type = doc_type
        self.content = []

    def add_content(self, content: Any):
        self.content.append(content)

    def __str__(self):
        return f"{self.name} ({self.type}):\n" + "\n".join(map(str, self.content))

class TextContent:
    def __init__(self, text: str):
        self.text = text

    def __str__(self):
        return self.text

class ChartContent:
    def __init__(self, chart_type: str, data: List[Dict[str, Any]]):
        self.chart_type = chart_type
        self.data = data

    def __str__(self):
        return f"[Chart: {self.chart_type}]\nData: {self.data}"


class TableContent:
    def __init__(self, headers: List[str], rows: List[List[Any]]):
        self.headers = headers
        self.rows = self._normalize_rows(rows, len(headers))

    def _normalize_rows(self, rows: List[List[Any]], header_count: int) -> List[List[str]]:
        return [
            [str(cell) if cell else "N/A" for cell in row[:header_count]] + ["N/A"] * (header_count - len(row))
            for row in rows
        ]

    def __str__(self):
        table_str = "| " + " | ".join(self.headers) + " |\n"
        table_str += "|" + "|".join(["---" for _ in self.headers]) + "|\n"
        for row in self.rows:
            table_str += "| " + " | ".join(row) + " |\n"
        return table_str

class FormulaContent:
    def __init__(self, formula: str, cell: str):
        self.formula = formula
        self.cell = cell

    def __str__(self):
        return f"Formula in {self.cell}: {self.formula}"

class ImageContent:
    def __init__(self, description: str):
        self.description = description

    def __str__(self):
        return f"[Image: {self.description}]"

class SlideContent:
    def __init__(self, number: int, content: str):
        self.number = number
        self.content = content

    def __str__(self):
        return f"--- Slide {self.number} ---\n{self.content}"
class DocumentGenerator:
    @staticmethod
    def generate_word_document(doc: 'CustomDocument') -> WordDocument:
        word_doc = WordDocument()

        # Set the default paragraph direction to RTL
        word_doc.styles['Normal'].paragraph_format.rtl = True

        # Add styles only if they don't exist
        styles = word_doc.styles
        if 'Heading 1' not in styles:
            style = styles.add_style('Heading 1', WD_STYLE_TYPE.PARAGRAPH)
            style.font.name = 'Arial'
            style.font.size = Pt(18)
            style.font.color.rgb = RGBColor(0, 0, 128)
            style.paragraph_format.rtl = True

        if 'Heading 2' not in styles:
            style = styles.add_style('Heading 2', WD_STYLE_TYPE.PARAGRAPH)
            style.font.name = 'Arial'
            style.font.size = Pt(16)
            style.font.color.rgb = RGBColor(0, 0, 100)
            style.paragraph_format.rtl = True

        if 'Normal' not in styles:
            style = styles.add_style('Normal', WD_STYLE_TYPE.PARAGRAPH)
            style.font.name = 'Arial'
            style.font.size = Pt(11)
            style.paragraph_format.rtl = True

        # Add title
        word_doc.add_paragraph(doc.name, style='Heading 1').alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        # Add introduction
        word_doc.add_paragraph("מסמך זה מכיל סיכום של הנתונים העסקיים העיקריים.",
                               style='Normal').alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT

        for item in doc.content:
            if isinstance(item, TextContent):
                p = word_doc.add_paragraph(style='Normal')
                run = p.add_run(item.text)
                p.alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT
            elif isinstance(item, ChartContent):
                word_doc.add_paragraph(f"תרשים: {item.chart_type}",
                                       style='Heading 2').alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT

                if item.chart_type == 'עוגה':
                    img_stream = DocumentGenerator._create_pie_chart(item.data)
                    try:
                        word_doc.add_picture(img_stream, width=Inches(6))
                        word_doc.paragraphs[-1].alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                    except Exception as e:
                        print(f"Error adding pie chart: {str(e)}")
                        word_doc.add_paragraph("שגיאה בהוספת תרשים העוגה", style='Normal')

                # Add data table for the chart
                table = word_doc.add_table(rows=len(item.data) + 1, cols=2)
                table.style = 'Table Grid'
                DocumentGenerator._set_cell_direction(table.cell(0, 0), 'RTL')
                DocumentGenerator._set_cell_direction(table.cell(0, 1), 'RTL')
                table.cell(0, 0).text = "קטגוריה"
                table.cell(0, 1).text = "ערך"
                for i, d in enumerate(item.data, start=1):
                    DocumentGenerator._set_cell_direction(table.cell(i, 0), 'RTL')
                    DocumentGenerator._set_cell_direction(table.cell(i, 1), 'RTL')
                    table.cell(i, 0).text = str(d['x'])
                    table.cell(i, 1).text = f"{d['y']:,}"

            elif isinstance(item, TableContent):
                word_doc.add_paragraph("נתונים טבלאיים", style='Heading 2').alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT
                table = word_doc.add_table(rows=len(item.rows) + 1, cols=len(item.headers))
                table.style = 'Table Grid'
                table.alignment = WD_TABLE_ALIGNMENT.RIGHT
                for i, header in enumerate(item.headers):
                    cell = table.cell(0, i)
                    DocumentGenerator._set_cell_direction(cell, 'RTL')
                    cell.text = header
                    cell.paragraphs[0].runs[0].font.bold = True
                for i, row in enumerate(item.rows, start=1):
                    for j, cell_value in enumerate(row):
                        cell = table.cell(i, j)
                        DocumentGenerator._set_cell_direction(cell, 'RTL')
                        cell.text = str(cell_value)



        return word_doc

    @staticmethod
    def _create_pie_chart(data):
        plt.figure(figsize=(8, 6))
        labels = [str(d['x']) for d in data]
        sizes = [d['y'] for d in data]
        colors = plt.cm.Set3(np.linspace(0, 1, len(sizes)))

        plt.pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%', startangle=90, textprops={'fontsize': 12})
        plt.axis('equal')
        plt.title("תרשים עוגה", fontsize=16, fontweight='bold')

        plt.rcParams['font.family'] = 'Arial'
        plt.rcParams['font.size'] = 12
        plt.rcParams['axes.unicode_minus'] = False

        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=300, facecolor='white', edgecolor='none')
        img_stream.seek(0)
        plt.close()

        return img_stream

    @staticmethod
    def _set_cell_direction(cell, direction):
        paragraph = cell.paragraphs[0]
        run = paragraph.runs
        for run in paragraph.runs:
            run.font.rtl = direction == 'RTL'

    @staticmethod
    def generate_spreadsheet(doc: 'CustomDocument') -> Workbook:
        wb = Workbook()
        ws = wb.active
        ws.title = doc.name
        ws.sheet_view.rightToLeft = True  # Right-to-left for Hebrew

        for item in doc.content:
            if isinstance(item, TextContent):
                ws.append([item.text])
                ws['A' + str(ws.max_row)].font = Font(bold=True, size=14)
            elif isinstance(item, ChartContent):
                ws.append([f"תרשים: {item.chart_type}"])
                ws.append(["קטגוריה", "ערך"])
                for d in item.data:
                    ws.append([d['x'], d['y']])

                chart = None
                if item.chart_type == 'עוגה':
                    chart = PieChart()
                elif item.chart_type == 'קווי':
                    chart = LineChart()
                else:
                    chart = BarChart()

                chart.title = item.chart_type
                data = Reference(ws, min_col=2, min_row=ws.max_row - len(item.data), max_row=ws.max_row)
                cats = Reference(ws, min_col=1, min_row=ws.max_row - len(item.data) + 1, max_row=ws.max_row)
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                ws.add_chart(chart, f"D{ws.max_row - len(item.data)}")
            elif isinstance(item, FormulaContent):
                ws[item.cell] = f"={item.formula}"
                red_text = Font(color="FF0000")
                red_fill = PatternFill(start_color="FFCCCC", end_color="FFCCCC", fill_type="solid")
                dxf = DifferentialStyle(font=red_text, fill=red_fill)
                rule = Rule(type="cellIs", operator="lessThan", formula=["0"], dxf=dxf)
                ws.conditional_formatting.add(item.cell, rule)
            elif isinstance(item, TableContent):
                ws.append(item.headers)
                for row in item.rows:
                    ws.append(row)

        for column_cells in ws.columns:
            length = max(len(str(cell.value)) for cell in column_cells)
            ws.column_dimensions[column_cells[0].column_letter].width = length + 2

        return wb

    @staticmethod
    def generate_presentation(doc: 'CustomDocument') -> Presentation:
        prs = Presentation()
        layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        if title:
            title.text = doc.name
            title.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(0, 0, 128)  # Dark blue
        if subtitle:
            subtitle.text = "נוצר עם מערכת ODSL משופרת"

        for item in doc.content:
            if isinstance(item, SlideContent):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                if slide.shapes.title:
                    slide.shapes.title.text = f"שקופית {item.number}"
                    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(0, 0, 128)  # Dark blue
                if len(slide.placeholders) > 1:
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.text = item.content
                    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT  # Right-align for Hebrew
                    for paragraph in tf.paragraphs:
                        paragraph.font.name = 'Arial'
                        paragraph.font.size = PptxPt(18)
            elif isinstance(item, ChartContent):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = f"תרשים: {item.chart_type}"
                chart_data = CategoryChartData()
                if isinstance(item.data, str):
                    item.data = eval(item.data)  # Convert string representation to actual list of dictionaries
                categories = [str(d['x']) for d in item.data]
                values = [d['y'] for d in item.data]
                # Check the type of item.data and handle accordingly
                if isinstance(item.data, list):
                    if all(isinstance(d, dict) for d in item.data):
                        categories = [str(d.get('x', '')) for d in item.data]
                        values = [d.get('y', 0) for d in item.data]
                    else:
                        categories = [str(i) for i in range(len(item.data))]
                        values = item.data
                elif isinstance(item.data, dict):
                    categories = [str(item.data.get('x', ''))]
                    values = [item.data.get('y', 0)]
                elif isinstance(item.data, str):
                    # If data is a string, use it as a single category
                    categories = [item.data]
                    values = [1]  # Default value
                else:
                    categories = ['No Data']
                    values = [0]

                chart_data.categories = categories
                chart_data.add_series('Series 1', values)

                x, y, cx, cy = PptxInches(2), PptxInches(2), PptxInches(6), PptxInches(4.5)
                chart_type = XL_CHART_TYPE.PIE if item.chart_type == 'עוגה' else XL_CHART_TYPE.LINE if item.chart_type == 'קווי' else XL_CHART_TYPE.COLUMN_CLUSTERED
                chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.has_title = True
                chart.chart_title.text_frame.text = item.chart_type
            elif isinstance(item, ImageContent):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                if slide.shapes.title:
                    slide.shapes.title.text = f"תמונה: {item.description}"
                if len(slide.placeholders) > 1:
                    body_shape = slide.placeholders[1]
                    body_shape.text = f"[מקום שמור לתמונה: {item.description}]"
        return prs

    @staticmethod
    def _create_chart_image(chart_content: ChartContent) -> io.BytesIO:
        fig, ax = plt.subplots()
        x = [d['x'] for d in chart_content.data]
        y = [d['y'] for d in chart_content.data]
        if chart_content.chart_type == 'עוגה':
            ax.pie(y, labels=x, autopct='%1.1f%%')
        elif chart_content.chart_type == 'עמודות':
            ax.bar(x, y)
        else:  # Default to line chart
            ax.plot(x, y)
        ax.set_title(chart_content.chart_type)
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png')
        img_stream.seek(0)
        plt.close(fig)
        return img_stream

    @staticmethod
    def _create_excel_chart(chart_content: ChartContent):
        chart = BarChart()
        chart.title = chart_content.chart_type
        chart.x_axis.title = 'X'
        chart.y_axis.title = 'Y'
        data = Reference(range_string=f"Sheet1!$A$1:$B${len(chart_content.data) + 1}")
        chart.add_data(data, titles_from_data=True)
        return chart

class MockOpenAI:
    def __init__(self, api_key=None):
        self.api_key = api_key
        self.chat = self.Chat()

    class Chat:
        def create(self, model: str, messages: List[Dict[str, str]], **kwargs) -> Any:
            user_input = messages[-1]['content']
            return MockOpenAI.ChatCompletion(self.generate_odsl_commands(user_input))

        def generate_odsl_commands(self, user_input: str) -> str:
            commands = []

            # Determine document type
            doc_type = "DOCUMENT"
            if "מצגת" in user_input or "presentation" in user_input.lower():
                doc_type = "PRESENTATION"
            elif "גיליון" in user_input or "spreadsheet" in user_input.lower():
                doc_type = "SPREADSHEET"

            # Extract document name
            doc_match = re.search(r"(?:מסמך|מצגת|גיליון)\s+['\"]([^'\"]+)['\"]", user_input)
            doc_name = doc_match.group(1) if doc_match else "UntitledDoc"

            commands.append(f"CREATE({doc_type}.{doc_name}, name='{doc_name}')")

            # Handle multiple slides for presentations
            if doc_type == "PRESENTATION":
                slide_contents = re.findall(r"שקופית\s+['\"]([^'\"]+)['\"]", user_input)
                for i, content in enumerate(slide_contents, 1):
                    commands.append(f"ADD_SLIDE({doc_type}.{doc_name}, number={i}, content='{content}')")
            else:
                # Extract title and content for non-presentation documents
                title_match = re.search(r"כותרת\s+['\"]([^'\"]+)['\"]", user_input)
                if title_match:
                    title = title_match.group(1)
                    commands.append(f"ADD({doc_type}.{doc_name}, content='{title}')")
                    commands.append(f"FORMAT({doc_type}.{doc_name}.LAST_PARAGRAPH, style='TITLE')")

                content_match = re.search(r"תוכן\s+['\"]([^'\"]+)['\"]", user_input)
                if content_match:
                    content = content_match.group(1)
                    commands.append(f"ADD({doc_type}.{doc_name}, content='{content}')")

            # Handle advanced formatting
            format_match = re.search(r"עיצוב\s+['\"]([^'\"]+)['\"]", user_input)
            if format_match:
                style = format_match.group(1)
                commands.append(f"FORMAT_ADVANCED({doc_type}.{doc_name}, style='{style}')")

            # Handle charts
            chart_match = re.search(r"תרשים\s+['\"]([^'\"]+)['\"]", user_input)
            if chart_match:
                chart_type = chart_match.group(1)
                commands.append(f"ADD_CHART({doc_type}.{doc_name}, type='{chart_type}')")

            # Handle images
            image_match = re.search(r"תמונה\s+['\"]([^'\"]+)['\"]", user_input)
            if image_match:
                image_desc = image_match.group(1)
                commands.append(f"ADD_IMAGE({doc_type}.{doc_name}, description='{image_desc}')")

            # Handle links
            link_match = re.search(r"קישור\s+['\"]([^'\"]+)['\"]", user_input)
            if link_match:
                link_url = link_match.group(1)
                commands.append(f"ADD_LINK({doc_type}.{doc_name}, url='{link_url}')")

            return "\n".join(commands)

    class ChatCompletion:
        def __init__(self, content):
            self.choices = [self.Choice(content)]

        class Choice:
            def __init__(self, content):
                self.message = self.Message(content)

            class Message:
                def __init__(self, content):
                    self.content = content




class ODSLCommand:
    def __init__(self, action: str, target: str, parameters: Dict[str, Any]):
        self.action = action
        self.target = target
        self.parameters = parameters

    def __str__(self):
        return f"ODSLCommand(action={self.action}, target={self.target}, parameters={self.parameters})"

    def to_dict(self):
        return {
            "action": self.action,
            "target": self.target,
            "parameters": self.parameters
        }

    @classmethod
    def from_dict(cls, data):
        return cls(data["action"], data["target"], data["parameters"])

class ODSLParser:
    def __init__(self):
        self.command_pattern = r'(\w+)\(([^,]+),\s*(.+)\)'

    def parse(self, odsl_input: str) -> List[ODSLCommand]:
        commands = []
        for line in odsl_input.strip().split('\n'):
            match = re.match(self.command_pattern, line.strip())
            if match:
                action, target, params_str = match.groups()
                parameters = self._parse_parameters(params_str)
                commands.append(ODSLCommand(action, target, parameters))
        return commands

    def _parse_parameters(self, params_str: str) -> Dict[str, Any]:
        params = {}
        pattern = r'(\w+)\s*=\s*(?:\'([^\']*?)\'|"([^"]*?)"|\[([^\]]*)\]|(\S+))'
        matches = re.findall(pattern, params_str)
        for match in matches:
            key, value1, value2, list_value, plain_value = match
            if list_value:
                try:
                    # Use ast.literal_eval for safer evaluation
                    params[key] = ast.literal_eval(f"[{list_value}]")
                except (SyntaxError, ValueError):
                    # If evaluation fails, keep it as a string
                    params[key] = list_value
            else:
                params[key] = value1 or value2 or plain_value
        return params








class ExecutionEngine:
    def __init__(self):
        self.documents = {}
        self.supported_actions = {
            "CREATE": self._execute_create,
            "ADD": self._execute_add,
            "FORMAT": self._execute_format,
            "ADD_SLIDE": self._execute_add_slide,
            "ADD_CHART": self._execute_add_chart,
            "ADD_IMAGE": self._execute_add_image,
            "ADD_LINK": self._execute_add_link,
            "FORMAT_ADVANCED": self._execute_format_advanced,
            "ADD_FORMULA": self._execute_add_formula,
            "ADD_TABLE": self._execute_add_table,
        }

    def execute(self, commands: List['ODSLCommand']):
        results = []
        for command in commands:
            result = self._execute_command(command)
            results.append(result)
        return results

    def _execute_command(self, command: 'ODSLCommand') -> str:
        if command.action in self.supported_actions:
            return self.supported_actions[command.action](command)
        else:
            return f"Unsupported action: {command.action}"

    def _execute_create(self, command: 'ODSLCommand'):
        doc_type, doc_name = command.target.split('.')
        self.documents[doc_name] = CustomDocument(doc_name, doc_type)
        return f"Created {doc_type}: {doc_name}"

    def _execute_add(self, command: 'ODSLCommand'):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            content = TextContent(command.parameters["content"])
            self.documents[doc_name].add_content(content)
            return f"Added content to {command.target}"
        else:
            return f"Error: Document {doc_name} not found"

    def _execute_format(self, command: 'ODSLCommand'):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            style = command.parameters.get('style', '')
            self.documents[doc_name].add_content(TextContent(f"Formatted with style: {style}"))
            return f"Formatted {command.target} with style: {style}"
        else:
            return f"Error: Document {doc_name} not found"

    def _execute_add_slide(self, command: 'ODSLCommand'):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            # Clean the number value and convert to int
            slide_number_str = command.parameters.get('number', '1')
            slide_number = int(slide_number_str.rstrip(','))  # Remove trailing comma if present
            content = command.parameters.get('content', '')
            slide = SlideContent(slide_number, content)
            self.documents[doc_name].add_content(slide)
            return f"Added slide {slide_number} to {command.target}"
        else:
            return f"Error: Document {doc_name} not found"


    def _execute_add_chart(self, command: 'ODSLCommand'):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            chart_type = command.parameters.get('type', '')
            data = command.parameters.get('data', [{'x': 1, 'y': 1}, {'x': 2, 'y': 2}, {'x': 3, 'y': 3}])
            chart = ChartContent(chart_type, data)
            self.documents[doc_name].add_content(chart)
            return f"Added {chart_type} chart to {command.target}"
        else:
            return f"Error: Document {doc_name} not found"

    def _execute_add_image(self, command: 'ODSLCommand'):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            description = command.parameters.get('description', '')
            image = ImageContent(description)
            self.documents[doc_name].add_content(image)
            return f"Added image to {command.target}: {description}"
        else:
            return f"Error: Document {doc_name} not found"

    def _execute_add_link(self, command: 'ODSLCommand'):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            url = command.parameters.get('url', '')
            text = command.parameters.get('text', url)
            self.documents[doc_name].add_content(TextContent(f"Added Link: [{text}]({url})"))
            return f"Added link to {command.target}: {text} ({url})"
        else:
            return f"Error: Document {doc_name} not found"

    def _execute_format_advanced(self, command: 'ODSLCommand'):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            style = command.parameters.get('style', '')
            self.documents[doc_name].add_content(TextContent(f"Advanced Formatting: {style}"))
            return f"Applied advanced formatting to {command.target}: {style}"
        else:
            return f"Error: Document {doc_name} not found"

    def _execute_add_formula(self, command: 'ODSLCommand'):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            formula = command.parameters.get('formula', '')
            cell = command.parameters.get('cell', 'A1')
            formula_content = FormulaContent(formula, cell)
            self.documents[doc_name].add_content(formula_content)
            return f"Added formula to {command.target}: {formula} in cell {cell}"
        else:
            return f"Error: Document {doc_name} not found"

    def _execute_add_table(self, command: 'ODSLCommand'):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            headers = self._ensure_list(command.parameters.get('headers', []))
            rows = self._ensure_list(command.parameters.get('rows', []))
            table = TableContent(headers, rows)
            self.documents[doc_name].add_content(table)
            return f"Added table to {command.target}"
        else:
            return f"Error: Document {doc_name} not found"

    def _ensure_list(self, value: Any) -> List[Any]:
        if isinstance(value, str):
            return self._parse_list(value)
        elif isinstance(value, list):
            return value
        else:
            raise ValueError(f"Unexpected value type for table data: {type(value)}")

    def _parse_list(self, list_str: str) -> List[Any]:
        # Remove brackets and split by comma
        items = list_str.strip('[]').split(',')
        # Strip whitespace and quotes from each item
        return [item.strip().strip("'\"") for item in items]


class QualityChecker:
    @staticmethod
    def check_document(document: 'CustomDocument') -> List[str]:
        issues = []
        if not document.content:
            issues.append(f"Warning: {document.name} is empty")
        for item in document.content:
            if isinstance(item, ChartContent):
                issues.extend(QualityChecker._check_chart(item))
            elif isinstance(item, TableContent):
                issues.extend(QualityChecker._check_table(item))
            elif isinstance(item, FormulaContent):
                issues.extend(QualityChecker._check_formula(item))
        return issues

    @staticmethod
    def _check_chart(chart: ChartContent) -> List[str]:
        issues = []
        if not chart.data:
            issues.append("Chart has no data")
        elif len(chart.data) < 3:
            issues.append("Chart has insufficient data points")
        return issues

    @staticmethod
    def _check_table(table: TableContent) -> List[str]:
        issues = []
        if not table.headers:
            issues.append("Table has no headers")
        if not table.rows:
            issues.append("Table has no data rows")
        elif any(len(row) != len(table.headers) for row in table.rows):
            issues.append("Table has mismatched row lengths")
        return issues

    @staticmethod
    def _check_formula(formula: FormulaContent) -> List[str]:
        issues = []
        if not formula.formula:
            issues.append("Formula is empty")
        elif formula.formula.count('(') != formula.formula.count(')'):
            issues.append("Formula has mismatched parentheses")
        return issues

class JigsawEngine:
    def __init__(self):
        self.feedback_data = []
        self.improvement_suggestions = defaultdict(list)
        self.action_improvements = {
            "ADD_CHART": self._improve_chart,
            "ADD_TABLE": self._improve_table,
            "ADD_FORMULA": self._improve_formula,
        }

    def collect_feedback(self, command: 'ODSLCommand', execution_result: str, user_rating: int):
        self.feedback_data.append({
            "command": command.__dict__,
            "result": execution_result,
            "rating": user_rating
        })

    def analyze_feedback(self):
        self.improvement_suggestions.clear()
        for entry in self.feedback_data:
            if entry["rating"] < 4:  # Consider ratings below 4 as needing improvement
                command = ODSLCommand(**entry["command"])
                self.improvement_suggestions[command.action].append(entry)

    def get_improvement_suggestions(self) -> Dict[str, List[Dict]]:
        return dict(self.improvement_suggestions)

    def apply_improvements(self, documents: Dict[str, 'CustomDocument']):
        for action, entries in self.improvement_suggestions.items():
            if action in self.action_improvements:
                self.action_improvements[action](documents, entries)

    def _improve_chart(self, documents: Dict[str, 'CustomDocument'], entries: List[Dict]):
        for doc in documents.values():
            for item in doc.content:
                if isinstance(item, ChartContent):
                    chart_type = item.chart_type
                    if chart_type == 'עוגה':
                        item.data = self._generate_pie_chart_data(doc.name)
                    elif chart_type in ['קווי', 'עמודות']:
                        item.data = self._generate_line_bar_chart_data(doc.name)
        print("Improved chart data generation")

    def _improve_table(self, documents: Dict[str, 'CustomDocument'], entries: List[Dict]):
        for doc in documents.values():
            for item in doc.content:
                if isinstance(item, TableContent):
                    item.headers = self._clean_headers(item.headers)
                    item.rows = self._generate_realistic_data(item.headers, doc.name)
        print("Improved table structure and content")

    def _generate_pie_chart_data(self, doc_name: str):
        if "תקציב" in doc_name:
            categories = ['משכורות', 'תפעול', 'שיווק', 'מחקר ופיתוח']
        else:
            categories = ['מוצר A', 'מוצר B', 'מוצר C', 'מוצר D']
        return [{'x': cat, 'y': random.randint(10, 100)} for cat in categories]

    def _generate_line_bar_chart_data(self, doc_name: str):
        months = ['ינו', 'פבר', 'מרץ', 'אפר', 'מאי', 'יונ']
        if "תקציב" in doc_name:
            return [{'x': month, 'y': random.randint(50000, 100000)} for month in months]
        else:
            return [{'x': month, 'y': random.randint(1000, 10000)} for month in months]

    def _generate_realistic_data(self, headers: List[str], doc_name: str) -> List[List[str]]:
        rows = []
        for _ in range(3):  # Generate 3 rows of data
            row = []
            for header in headers:
                if "חודש" in header.lower():
                    row.append(random.choice(['ינואר', 'פברואר', 'מרץ', 'אפריל', 'מאי', 'יוני']))
                elif "הכנסות" in header.lower():
                    row.append(str(random.randint(50000, 100000)))
                elif "הוצאות" in header.lower():
                    row.append(str(random.randint(40000, 90000)))
                else:
                    row.append(f"נתון {random.randint(1, 100)}")
            rows.append(row)
        return rows

    def _improve_formula(self, documents: Dict[str, 'CustomDocument'], entries: List[Dict]):
        for doc in documents.values():
            for item in doc.content:
                if isinstance(item, FormulaContent):
                    item.formula = self._optimize_formula(item.formula)
        print("Optimized spreadsheet formulas")

    def _clean_headers(self, headers):
        return [header.replace('Improved ', '').replace('Better ', '') for header in headers]

    def _clean_rows(self, rows):
        return [[cell.replace('Better ', '').replace('[', '').replace(']', '').replace("'", '') for cell in row] for row in rows]

    def _optimize_formula(self, old_formula):
        if "SUM" in old_formula:
            return old_formula.replace("SUM", "SUMIF")
        return old_formula


def get_openai_client(use_mock=True, api_key=None):
    if use_mock:
        return MockOpenAI(api_key)
    else:
        from openai import OpenAI
        return OpenAI(api_key=api_key)

class ChatGPTInterface:
    def __init__(self, api_key: str, use_mock: bool = True):
        self.client = MockOpenAI(api_key=api_key) if use_mock else OpenAI(api_key=api_key)

    def natural_language_to_odsl(self, user_input: str) -> str:
        response = self.client.chat.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": user_input}]
        )
        return response.choices[0].message.content


class MockNLToODSL:
    def convert(self, user_input: str) -> str:
        commands = []

        doc_match = re.search(r"(מצגת|מסמך|גיליון)\s+['\"](.*?)['\"]", user_input)
        if doc_match:
            doc_type = doc_match.group(1)
            doc_name = doc_match.group(2)

            if doc_type == "מצגת":
                doc_type_en = "PRESENTATION"
            elif doc_type == "מסמך":
                doc_type_en = "DOCUMENT"
            elif doc_type == "גיליון":
                doc_type_en = "SPREADSHEET"

            commands.append(f"CREATE({doc_type_en}.{doc_name}, name='{doc_name}')")

        slide_matches = re.findall(r"שקופית\s+['\"](.*?)['\"]", user_input)
        for i, slide_content in enumerate(slide_matches, 1):
            commands.append(f"ADD_SLIDE({doc_type_en}.{doc_name}, number={i}, content='{slide_content}')")

        content_match = re.search(r"תוכן\s+['\"](.*?)['\"]", user_input)
        if content_match:
            commands.append(f"ADD({doc_type_en}.{doc_name}, content='{content_match.group(1)}')")

        style_match = re.search(r"עיצוב\s+['\"](.*?)['\"]", user_input)
        if style_match:
            commands.append(f"FORMAT_ADVANCED({doc_type_en}.{doc_name}, style='{style_match.group(1)}')")

        chart_match = re.search(r"תרשים\s+['\"](.*?)['\"]", user_input)
        if chart_match:
            chart_type = chart_match.group(1)
            data = self._generate_chart_data(chart_type)
            commands.append(f"ADD_CHART({doc_type_en}.{doc_name}, type='{chart_type}', data={data})")

        formula_match = re.search(r"נוסחה\s+['\"](.*?)['\"]", user_input)
        if formula_match:
            formula = formula_match.group(1)
            cell = self._get_appropriate_cell(formula)
            commands.append(f"ADD_FORMULA({doc_type_en}.{doc_name}, formula='{formula}', cell='{cell}')")

        image_match = re.search(r"תמונה\s+['\"](.*?)['\"]", user_input)
        if image_match:
            commands.append(f"ADD_IMAGE({doc_type_en}.{doc_name}, description='{image_match.group(1)}')")

        table_match = re.search(r"טבלה\s+['\"](.*?)['\"]", user_input)
        if table_match:
            table_name = table_match.group(1)
            headers, rows = self._generate_table_data(table_name)
            commands.append(f"ADD_TABLE({doc_type_en}.{doc_name}, headers={headers}, rows={rows})")

        return "\n".join(commands)

    def _generate_chart_data(self, chart_type: str) -> str:
        if chart_type == "עוגה":
            return "[{'x': 'A', 'y': 30}, {'x': 'B', 'y': 50}, {'x': 'C', 'y': 20}]"
        else:
            return "[{'x': 1, 'y': 10}, {'x': 2, 'y': 20}, {'x': 3, 'y': 15}, {'x': 4, 'y': 25}, {'x': 5, 'y': 30}]"

    def _get_appropriate_cell(self, formula: str) -> str:
        if "SUM" in formula:
            return "A1"
        elif "AVERAGE" in formula:
            return "B1"
        else:
            return "C1"

    def _generate_table_data(self, table_name: str) -> tuple:
        if "כספי" in table_name:
            headers = "['חודש', 'הכנסות', 'הוצאות']"
            rows = "[['ינואר', '10000', '8000'], ['פברואר', '12000', '9000'], ['מרץ', '15000', '10000']]"
        else:
            headers = "['שם', 'גיל', 'עיר']"
            rows = "[['ישראל', '30', 'תל אביב'], ['שרה', '25', 'ירושלים'], ['יוסף', '35', 'חיפה']]"
        return headers, rows


class EnhancedExecutionEngine(ExecutionEngine):
    def __init__(self, use_mock_rating=False):
        super().__init__()
        self.jigsaw = JigsawEngine()
        self.use_mock_rating = use_mock_rating
        self.execution_history = []
        self.document_generator = DocumentGenerator()

    def execute(self, commands: List[ODSLCommand]):
        results = []
        for command in commands:
            try:
                result = self._execute_command(command)
                results.append(result)
            except Exception as e:
                print(f"Error executing command {command}: {str(e)}")
                print(f"Command details: {command.__dict__}")
        self._collect_feedback(commands, results)
        self.jigsaw.analyze_feedback()
        self._apply_improvements()
        self._check_quality()
        self._generate_real_documents()
        return results

    def _collect_feedback(self, commands: List[ODSLCommand], results: List[str]):
        for command, result in zip(commands, results):
            if self.use_mock_rating:
                user_rating = random.randint(1, 5)
            else:
                user_rating = int(input(f"Rate the execution of {command} (1-5): "))
            self.jigsaw.collect_feedback(command, result, user_rating)
            self.execution_history.append((command, result, user_rating))

    def _apply_improvements(self):
        suggestions = self.jigsaw.get_improvement_suggestions()
        for action, entries in suggestions.items():
            if action == 'CREATE':
                self._improve_create()
            elif action == 'ADD':
                self._improve_add()
            elif action == 'ADD_CHART':
                self._improve_charts()
            elif action == 'ADD_FORMULA':
                self._improve_formulas()

    def _improve_formulas(self):
        for doc in self.documents.values():
            for item in doc.content:
                if isinstance(item, FormulaContent):
                    if item.formula.startswith('SUM'):
                        item.formula = item.formula.replace('SUM', 'SUMIF')
                    item.cell = 'B1'  # Move formula to a more prominent position


    def _improve_formatting(self):
        for doc in self.documents.values():
            if doc.type == 'SPREADSHEET':
                for item in doc.content:
                    if isinstance(item, TextContent):
                        item.text = item.text.upper()  # Make text uppercase
                    elif isinstance(item, FormulaContent):
                        # Enhance formula if possible
                        if item.formula.startswith('SUM'):
                            item.formula = item.formula.replace('SUM', 'SUMIF')

    def _execute_add_chart(self, command: ODSLCommand):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            chart_type = command.parameters.get('type', '')
            data = command.parameters.get('data', [])

            # If no data is provided or data is empty, generate some
            if not data or (isinstance(data, list) and len(data) == 0):
                data = self._generate_chart_data(chart_type)

            chart = ChartContent(chart_type, data)
            self.documents[doc_name].add_content(chart)
            return f"Added {chart_type} chart to {command.target}"
        else:
            return f"Error: Document {doc_name} not found"

    def _generate_chart_data(self, chart_type):
        if chart_type == 'עוגה':
            categories = ['קטגוריה A', 'קטגוריה B', 'קטגוריה C', 'קטגוריה D', 'קטגוריה E']
            return [{'x': cat, 'y': random.randint(10, 100)} for cat in categories]
        else:
            return [{'x': i, 'y': random.randint(10, 100)} for i in range(1, 6)]


    def _improve_charts(self):
        for doc in self.documents.values():
            for item in doc.content:
                if isinstance(item, ChartContent):
                    if item.chart_type == 'עוגה':
                        item.data = [
                            {'x': 'מוצר א', 'y': 30},
                            {'x': 'מוצר ב', 'y': 25},
                            {'x': 'מוצר ג', 'y': 20},
                            {'x': 'מוצר ד', 'y': 15},
                            {'x': 'מוצר ה', 'y': 10}
                        ]
                    elif item.chart_type == 'קווי':
                        item.data = [
                            {'x': 'ינואר', 'y': 50},
                            {'x': 'פברואר', 'y': 60},
                            {'x': 'מרץ', 'y': 70},
                            {'x': 'אפריל', 'y': 65},
                            {'x': 'מאי', 'y': 80}
                        ]

    def _improve_slides(self):
        for doc in self.documents.values():
            if doc.type == 'PRESENTATION':
                for item in doc.content:
                    if isinstance(item, SlideContent):
                        item.content += "\n\nImproved content for better presentation."

    def _improve_images(self):
        for doc in self.documents.values():
            for item in doc.content:
                if isinstance(item, ImageContent):
                    item.description += " (High Resolution)"
    def _check_quality(self):
        issues = []
        for doc_name, doc in self.documents.items():
            doc_issues = QualityChecker.check_document(doc)
            if doc_issues:
                issues.extend([f"{doc_name}: {issue}" for issue in doc_issues])
        if issues:
            print("Quality issues found:")
            for issue in issues:
                print(f"  - {issue}")
        else:
            print("No quality issues found.")

    def _improve_create(self):
        for doc in self.documents.values():
            doc.content.insert(0, TextContent(f"מסמך מתקדם: {doc.name} - נוצר עם מערכת ODSL משופרת"))



    def _improve_add(self):
        for doc in self.documents.values():
            for item in doc.content:
                if isinstance(item, TextContent):
                    item.text = f"{item.text} - תוכן משופר"

    def _generate_real_documents(self):
        generated_files = []
        for doc_name, doc in self.documents.items():
            try:
                if doc.type == "DOCUMENT":
                    word_doc = self.document_generator.generate_word_document(doc)
                    filename = self._get_safe_filename(f"{doc_name}.docx")
                    word_doc.save(filename)
                    generated_files.append(filename)
                elif doc.type == "SPREADSHEET":
                    excel_doc = self.document_generator.generate_spreadsheet(doc)
                    filename = self._get_safe_filename(f"{doc_name}.xlsx")
                    excel_doc.save(filename)
                    generated_files.append(filename)
                elif doc.type == "PRESENTATION":
                    ppt_doc = self.document_generator.generate_presentation(doc)
                    filename = self._get_safe_filename(f"{doc_name}.pptx")
                    ppt_doc.save(filename)
                    generated_files.append(filename)
                print(f"Generated {doc.type} document: {filename}")
            except Exception as e:
                print(f"Error generating {doc.type} document '{doc_name}': {str(e)}")
                print(f"Document content:")
                for item in doc.content:
                    print(f"  - Type: {type(item)}")
                    if isinstance(item, ChartContent):
                        print(f"    Chart type: {item.chart_type}")
                        print(f"    Data type: {type(item.data)}")
                        print(f"    Data: {item.data}")
                    elif isinstance(item, TableContent):
                        print(f"    Headers: {item.headers}")
                        print(f"    Rows: {item.rows}")
                import traceback
                traceback.print_exc()
        return generated_files

    def _get_safe_filename(self, filename):
        base, ext = os.path.splitext(filename)
        counter = 1
        while os.path.exists(filename):
            filename = f"{base}_{counter}{ext}"
            counter += 1
        return filename


    def _safe_save(self, save_function, filename):
        try:
            save_function(filename)
        except PermissionError:
            print(f"Permission denied when saving {filename}. Trying with a modified name...")
            base, ext = os.path.splitext(filename)
            save_function(f"{base}_new{ext}")
        except Exception as e:
            print(f"Error saving {filename}: {str(e)}")

    def _improve_tables(self):
        for doc in self.documents.values():
            for item in doc.content:
                if isinstance(item, TableContent):
                    item.headers = ['חודש', 'הכנסות', 'הוצאות', 'רווח']
                    item.rows = [
                        ['ינואר', 50000, 30000, 20000],
                        ['פברואר', 55000, 32000, 23000],
                        ['מרץ', 60000, 35000, 25000],
                        ['אפריל', 58000, 33000, 25000],
                        ['מאי', 62000, 36000, 26000]
                    ]



    def _execute_command(self, command: ODSLCommand) -> str:
        if command.action in self.supported_actions:
            # Pre-process the command parameters if needed
            if command.action == 'ADD_CHART':
                command.parameters['data'] = self._preprocess_chart_data(command.parameters['data'])
            elif command.action == 'ADD_TABLE':
                command.parameters['headers'], command.parameters['rows'] = self._preprocess_table_data(
                    command.parameters['headers'], command.parameters['rows'])

            return self.supported_actions[command.action](command)
        else:
            return f"Unsupported action: {command.action}"

    def _preprocess_chart_data(self, data):
        if isinstance(data, str):
            try:
                return ast.literal_eval(data)
            except:
                print(f"Warning: Could not parse chart data: {data}")
                return data
        return data



    def _preprocess_table_data(self, headers, rows):
        if isinstance(headers, str):
            try:
                headers = ast.literal_eval(headers)
            except:
                print(f"Warning: Could not parse table headers: {headers}")
                headers = headers.strip("[]").replace("'", "").split(", ")

        if isinstance(rows, str):
            try:
                rows = ast.literal_eval(rows)
            except:
                print(f"Warning: Could not parse table rows: {rows}")
                # Try to parse the rows manually
                rows = rows.strip("[]").split("], [")
                rows = [row.replace("'", "").split(", ") for row in rows]

        # Ensure all rows have the same number of columns as headers
        max_cols = len(headers)
        rows = [row + [''] * (max_cols - len(row)) for row in rows]

        return headers, rows

    def _execute_add_table(self, command: ODSLCommand):
        doc_name = command.target.split('.')[-1]
        if doc_name in self.documents:
            headers = command.parameters.get('headers', [])
            rows = command.parameters.get('rows', [])

            # Always generate realistic data for financial tables
            headers, rows = self._generate_realistic_financial_data()

            table = TableContent(headers, rows)
            self.documents[doc_name].add_content(table)
            return f"Added table to {command.target}"
        else:
            return f"Error: Document {doc_name} not found"

    def _generate_realistic_financial_data(self):
        headers = ['חודש', 'הכנסות', 'הוצאות', 'רווח']
        months = ['ינואר', 'פברואר', 'מרץ', 'אפריל', 'מאי', 'יוני',
                  'יולי', 'אוגוסט', 'ספטמבר', 'אוקטובר', 'נובמבר', 'דצמבר']
        rows = []
        for month in months:
            income = random.randint(50000, 100000)
            expenses = random.randint(30000, 80000)
            profit = income - expenses
            row = [month, f"{income:,}", f"{expenses:,}", f"{profit:,}"]
            rows.append(row)
        return headers, rows

def test_system():
    parser = ODSLParser()
    execution_engine = EnhancedExecutionEngine(use_mock_rating=True)
    nl_to_odsl = MockNLToODSL()

    test_inputs = [
        "צור מסמך 'דוח שנתי' עם תוכן 'סיכום הישגים' ותרשים 'עוגה' וטבלה 'נתונים כספיים'",
        "הכן גיליון 'תקציב 2024' עם נוסחה 'SUM(A1:A10)' ועיצוב 'מקצועי' ותרשים 'קווי'",
        "הכן מצגת 'תכנית עסקית' עם שקופית 'חזון החברה' ושקופית 'יעדים' ותמונה 'לוגו' ותרשים 'עמודות'"
    ]

    for i, user_input in enumerate(test_inputs, 1):
        print(f"\n--- בדיקה {i} ---")
        print("קלט המשתמש:")
        print(user_input)
        odsl_commands = nl_to_odsl.convert(user_input)
        print("\nפקודות ODSL שנוצרו:")
        print(odsl_commands)


        commands = parser.parse(odsl_commands)
        execution_results = execution_engine.execute(commands)

        print("\nתוצאות הביצוע:")
        for result in execution_results:
            print(result)

    print("\n--- קבצים שנוצרו ---")
    generated_files = execution_engine._generate_real_documents()
    for file in generated_files:
        print(f"נוצר קובץ: {file}")

    print("\n--- סיכום ביצועי המערכת ---")
    print("מספר פקודות שבוצעו:", len(execution_engine.execution_history))
    print("פעולות שדורשות שיפור:", execution_engine.jigsaw.get_improvement_suggestions())

    quality_score = sum(1 for doc in execution_engine.documents.values() if not QualityChecker.check_document(doc))
    total_docs = len(execution_engine.documents)
    print(f"ציון איכות כולל: {quality_score}/{total_docs} מסמכים ללא בעיות")


if __name__ == "__main__":
    test_system()



















